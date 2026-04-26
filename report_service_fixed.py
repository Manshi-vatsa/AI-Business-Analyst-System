import logging
import io
import base64
from datetime import datetime
from typing import Dict, Any, List
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

class ReportService:
    """
    Service for generating PDF and PowerPoint reports
    """
    
    def __init__(self):
        pass
    
    def generate_pdf_report(self, summary: str, insights: List[str]) -> bytes:
        """
        Generate PDF report with summary and insights
        """
        try:
            # Create PDF buffer
            buffer = io.BytesIO()
            
            # Create PDF document
            doc = SimpleDocTemplate(buffer, pagesize=(8.5*inch, 11*inch), rightMargin=72)
            styles = getSampleStyleSheet()
            
            # Get style objects
            title_style = styles['Title']
            heading_style = styles['Heading2']
            normal_style = styles['Normal']
            
            # Build PDF content
            story = []
            
            # Title
            story.append(Paragraph("AI Analytics Report", title_style))
            story.append(Spacer(1, 12))
            
            # Report metadata
            story.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
            story.append(Spacer(1, 12))
            
            # Summary section
            story.append(Paragraph("Summary", heading_style))
            story.append(Spacer(1, 12))
            story.append(Paragraph(summary, normal_style))
            story.append(Spacer(1, 12))
            
            # Insights section
            story.append(Paragraph("Insights", heading_style))
            story.append(Spacer(1, 12))
            
            for insight in insights:
                story.append(Paragraph(f"• {insight}", normal_style))
                story.append(Spacer(1, 6))
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            logger.info("PDF report generated successfully")
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error generating PDF report: {e}")
            raise Exception(f"Failed to generate PDF report: {str(e)}")
    
    def generate_ppt_report(self, summary: str, insights: List[str]) -> bytes:
        """
        Generate PowerPoint report with summary and insights
        """
        try:
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = "AI Analytics Report"
            subtitle = title_slide.shapes.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # Summary slide
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            summary_title = summary_slide.shapes.title
            summary_title.text = "Summary"
            
            # Add summary content
            left = top = Inches(1)
            width, height = Inches(8), Inches(4)
            summary_textbox = summary_slide.shapes.add_textbox(left, top, width, height)
            summary_textbox.text_frame.text = summary
            
            # Insights slide
            insights_slide = prs.slides.add_slide(prs.slide_layouts[1])
            insights_title = insights_slide.shapes.title
            insights_title.text = "Insights"
            
            # Add insights content
            insights_text = "\n".join([f"• {insight}" for insight in insights])
            insights_textbox = insights_slide.shapes.add_textbox(left, top, width, height)
            insights_textbox.text_frame.text = insights_text
            
            # Save presentation to bytes
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            logger.info("PowerPoint report generated successfully")
            return ppt_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint report: {e}")
            raise Exception(f"Failed to generate PowerPoint report: {str(e)}")

# Global report service instance
report_service = ReportService()
