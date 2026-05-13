import logging
from typing import Dict, Any, Optional, List
from insights_service import insights_service
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches
import io
import base64
from datetime import datetime

logger = logging.getLogger(__name__)

class ReportService:
    """
    Service for generating PDF and PowerPoint reports
    """
    
    def __init__(self):
        pass
    
    def generate_pdf_report(self, summary: str, insights: List[str]) -> bytes:
        """
        Generate PDF report with summary and insights
        """
        try:
            # Create PDF buffer
            buffer = io.BytesIO()
            
            # Create PDF document
            doc = SimpleDocTemplate(buffer, pagesize=(8.5*inch, 11*inch), rightMargin=72)
            styles = getSampleStyleSheet()
            
            # Get style objects
            title_style = styles['Title']
            heading_style = styles['Heading2']
            normal_style = styles['Normal']
            
            # Build PDF content
            story = []
            
            # Title
            story.append(Paragraph("AI Analytics Report", title_style))
            story.append(Spacer(1, 12))
            
            # Report metadata
            story.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
            story.append(Spacer(1, 12))
            
            # Summary section
            story.append(Paragraph("Summary", heading_style))
            story.append(Spacer(1, 12))
            story.append(Paragraph(summary, normal_style))
            story.append(Spacer(1, 12))
            
            # Insights section
            story.append(Paragraph("Insights", heading_style))
            story.append(Spacer(1, 12))
            
            for insight in insights.get('insights', []):
                story.append(Paragraph(f"• {insight}", normal_style))
                story.append(Spacer(1, 6))
            
            # Regional Analysis
            if 'highest_region' in insights:
                story.append(Paragraph("Regional Analysis", heading_style))
                story.append(Spacer(1, 12))
                story.append(Paragraph(f"Highest Performing Region: {insights['highest_region']}"))
                story.append(Paragraph(f"Revenue: ${insights.get('highest_region_revenue', 0):,.2f}"))
                story.append(Spacer(1, 12))
            
            # Product Analysis
            if 'lowest_product' in insights:
                story.append(Paragraph("Product Analysis", heading_style))
                story.append(Spacer(1, 12))
                story.append(Paragraph(f"Lowest Performing Product: {insights['lowest_product']}"))
                story.append(Paragraph(f"Revenue: ${insights.get('lowest_product_revenue', 0):,.2f}"))
                story.append(Spacer(1, 12))
            
            # Performance Indicators
            if 'performance_indicators' in insights:
                story.append(Paragraph("Performance Indicators", heading_style))
                story.append(Spacer(1, 12))
                for indicator in insights['performance_indicators']:
                    story.append(Paragraph(f"• {indicator}", normal_style))
                    story.append(Spacer(1, 6))
            
            # Detailed Breakdowns
            if 'region_breakdown' in insights:
                story.append(Paragraph("Regional Breakdown", heading_style))
                story.append(Spacer(1, 12))
                for region in insights['region_breakdown']:
                    story.append(Paragraph(f"{region.get('region', 'N/A')}: ${region.get('revenue', 0):,.2f} ({region.get('count', 0)} sales)", normal_style))
                    story.append(Spacer(1, 6))
            
            if 'product_breakdown' in insights:
                story.append(Paragraph("Product Breakdown", heading_style))
                story.append(Spacer(1, 12))
                for product in insights['product_breakdown']:
                    story.append(Paragraph(f"{product.get('product', 'N/A')}: ${product.get('revenue', 0):,.2f} ({product.get('count', 0)} sales)", normal_style))
                    story.append(Spacer(1, 6))
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            logger.info("PDF report generated successfully")
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error generating PDF report: {e}")
            raise Exception(f"Failed to generate PDF report: {str(e)}")
    
    def generate_ppt_report(self, insights: Dict[str, Any]) -> bytes:
        """
        Generate PowerPoint report with summary and insights
        """
        try:
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = "AI Business Analytics Report"
            subtitle = title_slide.shapes.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # Executive Summary slide
            if 'total_revenue' in insights:
                summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
                summary_title = summary_slide.shapes.title
                summary_title.text = "Executive Summary"
                
                # Add summary content
                summary_content = f"""Total Revenue: ${insights.get('total_revenue', 0):,.2f}
Total Sales Count: {insights.get('total_sales_count', 0)}
Average Sale: ${insights.get('average_sale', 0):,.2f}"""
                
                summary_textbox = summary_slide.shapes.add_textbox(
                    summary_content, 
                    left=Inches(0.5), 
                    top=Inches(1.5), 
                    width=Inches(8), 
                    height=Inches(4)
                )
            
            # Regional Analysis slide
            if 'highest_region' in insights:
                region_slide = prs.slides.add_slide(prs.slide_layouts[1])
                region_title = region_slide.shapes.title
                region_title.text = "Regional Analysis"
                
                highest_region = insights['highest_region']
                region_content = f"""Highest Performing Region: {highest_region.get('region', 'N/A')}
Revenue: ${highest_region.get('revenue', 0):,.2f}"""
                
                region_textbox = region_slide.shapes.add_textbox(
                    region_content,
                    left=Inches(0.5),
                    top=Inches(1.5),
                    width=Inches(8),
                    height=Inches(4)
                )
            
            # Product Analysis slide
            if 'lowest_product' in insights:
                product_slide = prs.slides.add_slide(prs.slide_layouts[1])
                product_title = product_slide.shapes.title
                product_title.text = "Product Analysis"
                
                lowest_product = insights['lowest_product']
                product_content = f"""Lowest Performing Product: {lowest_product.get('product', 'N/A')}
Revenue: ${lowest_product.get('revenue', 0):,.2f}"""
                
                product_textbox = product_slide.shapes.add_textbox(
                    product_content,
                    left=Inches(0.5),
                    top=Inches(1.5),
                    width=Inches(8),
                    height=Inches(4)
                )
            
            # Performance Indicators slide
            if 'performance_indicators' in insights:
                performance_slide = prs.slides.add_slide(prs.slide_layouts[1])
                performance_title = performance_slide.shapes.title
                performance_title.text = "Performance Indicators"
                
                performance_content = "\n".join([f"• {indicator}" for indicator in insights['performance_indicators']])
                
                performance_textbox = performance_slide.shapes.add_textbox(
                    performance_content,
                    left=Inches(0.5),
                    top=Inches(1.5),
                    width=Inches(8),
                    height=Inches(4)
                )
            
            # Save presentation to bytes
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            logger.info("PowerPoint report generated successfully")
            return ppt_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint report: {e}")
            raise Exception(f"Failed to generate PowerPoint report: {str(e)}")
    
    def generate_report(self, report_type: str) -> Dict[str, Any]:
        """
        Generate report of specified type
        """
        try:
            # Get latest insights
            insights = insights_service.analyze_sales_data()
            
            if 'error' in insights:
                return {
                    "status": "error",
                    "error": insights['error'],
                    "message": "Failed to generate insights for report"
                }
            
            # Generate report based on type
            if report_type.lower() == 'pdf':
                report_data = self.generate_pdf_report(insights)
                content_type = 'application/pdf'
                filename = f'business_analytics_report_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pdf'
            elif report_type.lower() == 'ppt':
                report_data = self.generate_ppt_report(insights)
                content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                filename = f'business_analytics_report_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
            else:
                return {
                    "status": "error",
                    "error": f"Unsupported report type: {report_type}",
                    "message": "Supported types: pdf, ppt"
                }
            
            # Encode report data for transmission
            encoded_data = base64.b64encode(report_data).decode('utf-8')
            
            return {
                "status": "success",
                "report_type": report_type,
                "filename": filename,
                "content_type": content_type,
                "data": encoded_data,
                "insights": insights,
                "message": f"{report_type.upper()} report generated successfully"
            }
            
        except Exception as e:
            logger.error(f"Error generating report: {e}")
            return {
                "status": "error",
                "error": str(e),
                "message": f"Failed to generate {report_type} report"
            }

# Global report service instance
report_service = ReportService()
